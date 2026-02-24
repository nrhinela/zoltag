"""Router for photo list operations."""

import io
import logging
import os
import re
import tempfile

from datetime import datetime
from typing import Optional
from uuid import UUID, uuid4

from fastapi import APIRouter, Depends, HTTPException, Body, Query, UploadFile, File, Form
from fastapi.responses import FileResponse
from google.api_core.exceptions import NotFound
from google.cloud import storage
from PIL import Image, UnidentifiedImageError
from sqlalchemy import func, and_, or_
from sqlalchemy.orm import Session
from starlette.background import BackgroundTask
from starlette.concurrency import run_in_threadpool

from zoltag.asset_helpers import AssetReadinessError, load_assets_for_images, resolve_image_storage
from zoltag.dependencies import get_db, get_secret, get_tenant, get_tenant_setting
from zoltag.list_visibility import (
    can_edit_list,
    can_view_list,
    get_list_scope_clause,
    is_list_owner,
    is_tenant_admin_user,
    normalize_list_visibility,
    normalize_list_scope,
)
from zoltag.tenant import Tenant
from zoltag.models.config import PhotoList, PhotoListItem, Keyword, KeywordCategory
from zoltag.models.sharing import PresentationTemplate
from zoltag.metadata import AssetDerivative, ImageMetadata, MachineTag, Permatag
from zoltag.tagging import calculate_tags
from zoltag.models.requests import AddPhotoRequest, ReorderListItemsRequest
from zoltag.settings import settings
from zoltag.auth.dependencies import get_current_user
from zoltag.auth.models import UserProfile
from zoltag.routers.images._shared import _build_source_url, _resolve_provider_ref
from zoltag.storage import create_storage_provider
from zoltag.text_index import rebuild_asset_text_index
from zoltag.tenant_scope import assign_tenant_scope, tenant_column_filter, tenant_column_filter_for_values

router = APIRouter(
    prefix="/api/v1/lists",
    tags=["lists"]
)
logger = logging.getLogger(__name__)
PPTX_EXPORT_MAX_ITEMS = 120
PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPTX_SLIDE_WIDTH_IN = 13.333
PPTX_SLIDE_HEIGHT_IN = 7.5
PPTX_SLIDE_MARGIN_IN = 0.2
PPTX_TEMPLATE_MAX_UPLOAD_BYTES = 50 * 1024 * 1024
PPTX_TEMPLATE_ALLOWED_VISIBILITY = {"shared", "private"}


def _refresh_asset_text_index_for_asset(db: Session, tenant: Tenant, asset_id) -> None:
    if not asset_id:
        return
    try:
        rebuild_asset_text_index(
            db,
            tenant_id=tenant.id,
            asset_id=str(asset_id),
            include_embeddings=False,
        )
    except Exception as exc:  # noqa: BLE001
        logger.warning("Failed to refresh asset_text_index for asset %s: %s", asset_id, exc)


def _resolve_storage_or_409(
    *,
    image: ImageMetadata,
    tenant: Tenant,
    db: Session,
    assets_by_id=None,
    require_source: bool = False,
):
    try:
        return resolve_image_storage(
            image=image,
            tenant=tenant,
            db=None if assets_by_id is not None else db,
            assets_by_id=assets_by_id,
            strict=settings.asset_strict_reads,
            require_source=require_source,
        )
    except AssetReadinessError as exc:
        raise HTTPException(status_code=409, detail=str(exc))


def _cleanup_temp_file(path: str) -> None:
    try:
        if path:
            os.remove(path)
    except FileNotFoundError:
        return
    except Exception as exc:  # noqa: BLE001
        logger.warning("Failed to remove temporary export file %s: %s", path, exc)


def _sanitize_export_filename(title: str, suffix: str) -> str:
    base = str(title or "list").strip() or "list"
    base = re.sub(r"\s+", "-", base)
    base = re.sub(r"[^A-Za-z0-9._-]", "", base).strip("._-")
    if not base:
        base = "list"
    return f"{base}{suffix}"


def _normalize_template_visibility(value: Optional[str], default: str = "private") -> str:
    normalized = str(value or default).strip().lower()
    if normalized not in PPTX_TEMPLATE_ALLOWED_VISIBILITY:
        raise HTTPException(status_code=400, detail="visibility must be 'shared' or 'private'")
    return normalized


def _template_visible_to_user(
    *,
    template: PresentationTemplate,
    current_user: UserProfile,
    is_tenant_admin: bool,
) -> bool:
    if is_tenant_admin:
        return True
    if str(template.visibility or "").strip().lower() == "shared":
        return True
    return str(template.created_by_uid) == str(current_user.supabase_uid)


def _serialize_presentation_template(
    template: PresentationTemplate,
    *,
    owner_name: Optional[str] = None,
) -> dict:
    return {
        "id": str(template.id),
        "name": str(template.name or "").strip(),
        "visibility": str(template.visibility or "private").strip().lower(),
        "storage_key": str(template.storage_key or "").strip(),
        "original_filename": str(template.original_filename or "").strip(),
        "created_by_uid": str(template.created_by_uid),
        "created_by_name": owner_name or "",
        "created_at": template.created_at.isoformat() if template.created_at else None,
        "updated_at": template.updated_at.isoformat() if template.updated_at else None,
    }


def _prepare_image_bytes_for_pptx(file_bytes: bytes) -> tuple[bytes, tuple[int, int]]:
    """Return pptx-compatible image bytes and dimensions."""
    try:
        with Image.open(io.BytesIO(file_bytes)) as image:
            width, height = image.size
            image_format = (image.format or "").upper()
            if image_format in {"JPEG", "JPG", "PNG", "BMP", "GIF", "TIFF"}:
                if image_format in {"JPEG", "JPG"} and image.mode not in {"RGB", "L", "CMYK", "YCbCr"}:
                    converted = image.convert("RGB")
                    out = io.BytesIO()
                    converted.save(out, format="JPEG", quality=95, optimize=True)
                    return out.getvalue(), (width, height)
                return file_bytes, (width, height)

            if image.mode in {"RGBA", "LA"} or (image.mode == "P" and "transparency" in image.info):
                converted = image.convert("RGBA")
                out = io.BytesIO()
                converted.save(out, format="PNG", optimize=True)
                return out.getvalue(), (width, height)

            converted = image.convert("RGB")
            out = io.BytesIO()
            converted.save(out, format="JPEG", quality=95, optimize=True)
            return out.getvalue(), (width, height)
    except UnidentifiedImageError as exc:
        raise ValueError(f"Unsupported image format: {exc}") from exc


def _fit_image_to_slide(
    image_size: tuple[int, int],
    *,
    slide_width: int,
    slide_height: int,
    margin: int,
) -> tuple[int, int, int, int]:
    image_width, image_height = image_size
    if image_width <= 0 or image_height <= 0:
        raise ValueError("Invalid image dimensions")

    available_width = max(1, slide_width - (2 * margin))
    available_height = max(1, slide_height - (2 * margin))
    image_aspect = image_width / image_height
    available_aspect = available_width / available_height

    if image_aspect >= available_aspect:
        target_width = available_width
        target_height = max(1, int(target_width / image_aspect))
    else:
        target_height = available_height
        target_width = max(1, int(target_height * image_aspect))

    left = int((slide_width - target_width) / 2)
    top = int((slide_height - target_height) / 2)
    return left, top, target_width, target_height


def _build_list_pptx_file(
    *,
    export_rows: list[dict],
    tenant: Tenant,
    template_path: Optional[str] = None,
) -> tuple[str, int, int]:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("PPTX export requires python-pptx. Install dependency and redeploy.") from exc

    providers = {}
    presentation = Presentation(template_path) if template_path else Presentation()
    if not template_path:
        presentation.slide_width = Inches(PPTX_SLIDE_WIDTH_IN)
        presentation.slide_height = Inches(PPTX_SLIDE_HEIGHT_IN)
    blank_layout = min(
        list(presentation.slide_layouts),
        key=lambda layout: len(getattr(layout, "placeholders", [])),
    )
    margin = int(Inches(PPTX_SLIDE_MARGIN_IN))

    added = 0
    skipped = 0
    for row in export_rows:
        provider_name = str(row.get("provider_name") or "").strip().lower()
        source_ref = str(row.get("source_ref") or "").strip()
        if not provider_name or not source_ref:
            skipped += 1
            continue

        provider = providers.get(provider_name)
        if provider is None:
            provider = create_storage_provider(provider_name, tenant=tenant, get_secret=get_secret)
            providers[provider_name] = provider

        try:
            original_bytes = provider.download_file(source_ref)
            image_bytes, image_size = _prepare_image_bytes_for_pptx(original_bytes)
            left, top, width, height = _fit_image_to_slide(
                image_size,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
                margin=margin,
            )
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                left=left,
                top=top,
                width=width,
                height=height,
            )
            added += 1
        except Exception as exc:  # noqa: BLE001
            skipped += 1
            logger.warning(
                "Skipping list export image %s (%s): %s",
                row.get("image_id"),
                source_ref,
                exc,
            )

    if added < 1:
        raise RuntimeError("No exportable images were found for this list.")

    temp_file = tempfile.NamedTemporaryFile(prefix="zoltag-list-export-", suffix=".pptx", delete=False)
    temp_file.close()
    presentation.save(temp_file.name)
    return temp_file.name, added, skipped


def _tenant_filter(model, tenant: Tenant | str):
    if isinstance(tenant, Tenant):
        return tenant_column_filter(model, tenant)
    return tenant_column_filter_for_values(model, tenant, tenant)


def _list_permissions(*, list_row: PhotoList, current_user: UserProfile, is_tenant_admin: bool) -> dict:
    return {
        "visibility": normalize_list_visibility(getattr(list_row, "visibility", None)),
        "is_owner": is_list_owner(list_row, current_user),
        "can_edit": can_edit_list(list_row, user=current_user, is_tenant_admin=is_tenant_admin),
    }


def _serialize_list_row(
    *,
    list_row: PhotoList,
    created_by_name: str | None,
    item_count: int,
    current_user: UserProfile,
    is_tenant_admin: bool,
) -> dict:
    permissions = _list_permissions(
        list_row=list_row,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    return {
        "id": list_row.id,
        "title": list_row.title,
        "notebox": list_row.notebox,
        "visibility": permissions["visibility"],
        "is_owner": permissions["is_owner"],
        "can_edit": permissions["can_edit"],
        "created_at": list_row.created_at,
        "updated_at": list_row.updated_at,
        "created_by_uid": list_row.created_by_uid,
        "created_by_name": created_by_name,
        "item_count": int(item_count or 0),
    }


def _query_lists_for_scope(
    *,
    db: Session,
    tenant: Tenant | str,
    current_user: UserProfile,
    scope: str = "default",
    is_tenant_admin: bool,
):
    query = db.query(PhotoList).filter(_tenant_filter(PhotoList, tenant))
    clause = get_list_scope_clause(user=current_user, scope=scope, is_tenant_admin=is_tenant_admin)
    if clause is not None:
        query = query.filter(clause)
    return query


def get_most_recent_list(
    db: Session,
    tenant: Tenant | str,
    *,
    current_user: UserProfile,
    is_tenant_admin: bool,
):
    """Get the most recently created accessible list for a tenant."""
    return (
        _query_lists_for_scope(
            db=db,
            tenant=tenant,
            current_user=current_user,
            scope="default",
            is_tenant_admin=is_tenant_admin,
        )
        .order_by(PhotoList.created_at.desc(), PhotoList.id.desc())
        .first()
    )


def _get_accessible_list_or_404(
    *,
    db: Session,
    tenant: Tenant,
    list_id: int,
    current_user: UserProfile,
    is_tenant_admin: bool,
) -> PhotoList:
    list_row = db.query(PhotoList).filter(
        PhotoList.id == list_id,
        tenant_column_filter(PhotoList, tenant),
    ).first()
    if not list_row:
        raise HTTPException(status_code=404, detail="List not found")
    if not can_view_list(list_row, user=current_user, is_tenant_admin=is_tenant_admin):
        raise HTTPException(status_code=404, detail="List not found")
    return list_row


def _get_template_or_404(
    *,
    db: Session,
    tenant: Tenant,
    template_id: str,
) -> PresentationTemplate:
    try:
        parsed_id = UUID(str(template_id))
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="template_id must be a valid UUID")
    template = db.query(PresentationTemplate).filter(
        PresentationTemplate.id == parsed_id,
        tenant_column_filter(PresentationTemplate, tenant),
    ).first()
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    return template


@router.get("/presentation-templates", response_model=list)
async def list_presentation_templates(
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """List PPTX templates visible to the current user in this tenant."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    query = db.query(PresentationTemplate).filter(tenant_column_filter(PresentationTemplate, tenant))
    if not is_tenant_admin:
        query = query.filter(
            or_(
                PresentationTemplate.visibility == "shared",
                PresentationTemplate.created_by_uid == current_user.supabase_uid,
            )
        )
    rows = (
        query
        .order_by(
            func.lower(PresentationTemplate.name).asc(),
            PresentationTemplate.created_at.desc(),
        )
        .all()
    )
    owner_ids = sorted({row.created_by_uid for row in rows if row.created_by_uid})
    owner_name_map = {}
    if owner_ids:
        owner_name_map = {
            str(uid): (display_name or "")
            for uid, display_name in db.query(UserProfile.supabase_uid, UserProfile.display_name).filter(
                UserProfile.supabase_uid.in_(owner_ids)
            ).all()
        }
    return [
        _serialize_presentation_template(
            row,
            owner_name=owner_name_map.get(str(row.created_by_uid), ""),
        )
        for row in rows
    ]


@router.post("/presentation-templates", response_model=dict)
async def upload_presentation_template(
    file: UploadFile = File(...),
    name: Optional[str] = Form(None),
    visibility: str = Form("private"),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Upload a reusable PPTX template for this tenant."""
    normalized_visibility = _normalize_template_visibility(visibility)
    filename = str(file.filename or "").strip() or "template.pptx"
    if not filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported.")
    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")
    if len(file_bytes) > PPTX_TEMPLATE_MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=400,
            detail=f"Template file is too large (max {int(PPTX_TEMPLATE_MAX_UPLOAD_BYTES / (1024 * 1024))} MB).",
        )

    template_id = uuid4()
    storage_key = tenant.get_presentation_template_key(str(template_id), filename)
    resolved_name = str(name or "").strip() or filename.rsplit(".", 1)[0]
    template_row = assign_tenant_scope(
        PresentationTemplate(
            id=template_id,
            name=resolved_name,
            storage_key=storage_key,
            original_filename=filename,
            visibility=normalized_visibility,
            created_by_uid=current_user.supabase_uid,
        ),
        tenant,
    )

    bucket = storage.Client(project=settings.gcp_project_id).bucket(tenant.get_storage_bucket(settings))
    blob = bucket.blob(storage_key)
    try:
        blob.upload_from_string(
            file_bytes,
            content_type=file.content_type or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        db.add(template_row)
        db.commit()
        db.refresh(template_row)
    except Exception as exc:  # noqa: BLE001
        db.rollback()
        try:
            blob.delete()
        except Exception:
            pass
        raise HTTPException(status_code=500, detail=f"Failed to upload template: {exc}")

    return _serialize_presentation_template(
        template_row,
        owner_name=(current_user.display_name or ""),
    )


@router.patch("/presentation-templates/{template_id}", response_model=dict)
async def update_presentation_template(
    template_id: str,
    name: Optional[str] = Body(None),
    visibility: Optional[str] = Body(None),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Update a PPTX template metadata row."""
    template = _get_template_or_404(db=db, tenant=tenant, template_id=template_id)
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    if str(template.created_by_uid) != str(current_user.supabase_uid) and not is_tenant_admin:
        raise HTTPException(status_code=403, detail="Only template owner or admin can edit this template")

    if name is not None:
        resolved_name = str(name).strip()
        if not resolved_name:
            raise HTTPException(status_code=400, detail="name cannot be empty")
        template.name = resolved_name
    if visibility is not None:
        template.visibility = _normalize_template_visibility(visibility)

    db.commit()
    db.refresh(template)
    owner_name = (
        db.query(UserProfile.display_name).filter(UserProfile.supabase_uid == template.created_by_uid).scalar()
        or ""
    )
    return _serialize_presentation_template(template, owner_name=owner_name)


@router.delete("/presentation-templates/{template_id}", response_model=dict)
async def delete_presentation_template(
    template_id: str,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Delete a PPTX template and its stored object."""
    template = _get_template_or_404(db=db, tenant=tenant, template_id=template_id)
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    if str(template.created_by_uid) != str(current_user.supabase_uid) and not is_tenant_admin:
        raise HTTPException(status_code=403, detail="Only template owner or admin can delete this template")

    storage_key = str(template.storage_key or "").strip()
    db.delete(template)
    db.commit()

    if storage_key:
        try:
            bucket = storage.Client(project=settings.gcp_project_id).bucket(tenant.get_storage_bucket(settings))
            bucket.blob(storage_key).delete()
        except NotFound:
            pass
        except Exception:
            logger.warning("Failed to delete template object %s", storage_key, exc_info=True)

    return {"deleted": True, "template_id": template_id}


@router.get("/recent", response_model=dict)
async def get_recent_list(
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Get the most recently created list for the tenant (if any)."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    recent = get_most_recent_list(
        db,
        tenant,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if not recent:
        return {}
    created_by_name = None
    if recent.created_by_uid:
        user = db.query(UserProfile.display_name).filter(
            UserProfile.supabase_uid == recent.created_by_uid
        ).first()
        created_by_name = user[0] if user else None
    item_count = db.query(func.count(PhotoListItem.id)).filter(
        PhotoListItem.list_id == recent.id
    ).scalar() or 0
    return {
        **_serialize_list_row(
            list_row=recent,
            created_by_name=created_by_name,
            item_count=item_count,
            current_user=current_user,
            is_tenant_admin=is_tenant_admin,
        )
    }


@router.get("/{list_id:int}", response_model=dict)
async def get_list(
    list_id: int,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Get a single list by ID for the tenant."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    lst = _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    item_count = db.query(func.count(PhotoListItem.id)).filter(
        PhotoListItem.list_id == lst.id
    ).scalar() or 0
    created_by_name = None
    if lst.created_by_uid:
        user = db.query(UserProfile.display_name).filter(
            UserProfile.supabase_uid == lst.created_by_uid
        ).first()
        created_by_name = user[0] if user else None
    return _serialize_list_row(
        list_row=lst,
        created_by_name=created_by_name,
        item_count=item_count,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )


@router.delete("/items/{item_id}", response_model=dict)
async def delete_list_item(
    item_id: int,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Remove a photo from a list by PhotoListItem ID (must belong to tenant)."""
    item_row = db.query(PhotoListItem, PhotoList).join(
        PhotoList, PhotoListItem.list_id == PhotoList.id
    ).filter(
        PhotoListItem.id == item_id,
        tenant_column_filter(PhotoList, tenant),
    ).first()
    if not item_row:
        raise HTTPException(status_code=404, detail="List item not found")
    item, list_row = item_row
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    if not can_view_list(list_row, user=current_user, is_tenant_admin=is_tenant_admin):
        raise HTTPException(status_code=404, detail="List item not found")
    removed_asset_id = item.asset_id
    db.delete(item)
    db.commit()
    _refresh_asset_text_index_for_asset(db, tenant, removed_asset_id)
    return {"deleted": True, "item_id": item_id}


@router.post("", response_model=dict)
async def create_list(
    title: str = Body(...),
    notebox: Optional[str] = Body(None),
    visibility: Optional[str] = Body("shared"),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user)
):
    """Create a new list."""
    new_list = assign_tenant_scope(PhotoList(
        title=title,
        notebox=notebox,
        visibility=normalize_list_visibility(visibility, default="shared"),
        created_by_uid=current_user.supabase_uid if current_user else None
    ), tenant)
    db.add(new_list)
    db.commit()
    db.refresh(new_list)
    return {
        "id": new_list.id,
        "title": new_list.title,
        "visibility": normalize_list_visibility(new_list.visibility),
    }


@router.get("", response_model=list)
async def list_lists(
    visibility_scope: str = Query(default="default"),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """List tenant lists by visibility scope."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    normalized_scope = normalize_list_scope(visibility_scope)
    lists = (
        _query_lists_for_scope(
            db=db,
            tenant=tenant,
            current_user=current_user,
            scope=normalized_scope,
            is_tenant_admin=is_tenant_admin,
        )
        .order_by(func.lower(PhotoList.title).asc(), PhotoList.created_at.asc(), PhotoList.id.asc())
        .all()
    )
    counts = dict(
        db.query(PhotoListItem.list_id, func.count(PhotoListItem.id))
        .join(PhotoList, PhotoListItem.list_id == PhotoList.id)
        .filter(tenant_column_filter(PhotoList, tenant))
        .group_by(PhotoListItem.list_id)
        .all()
    )
    # Build a map of user UUIDs to display names
    user_names = {}
    if lists:
        uids = [l.created_by_uid for l in lists if l.created_by_uid]
        if uids:
            users = db.query(UserProfile.supabase_uid, UserProfile.display_name).filter(
                UserProfile.supabase_uid.in_(uids)
            ).all()
            user_names = {u[0]: u[1] for u in users}

    return [
        _serialize_list_row(
            list_row=l,
            created_by_name=user_names.get(l.created_by_uid) if l.created_by_uid else None,
            item_count=counts.get(l.id, 0),
            current_user=current_user,
            is_tenant_admin=is_tenant_admin,
        )
        for l in lists
    ]


@router.patch("/{list_id:int}", response_model=dict)
async def edit_list(
    list_id: int,
    title: Optional[str] = Body(None),
    notebox: Optional[str] = Body(None),
    visibility: Optional[str] = Body(None),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Edit a list (title and/or notebox)."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    lst = _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if not can_edit_list(lst, user=current_user, is_tenant_admin=is_tenant_admin):
        raise HTTPException(status_code=403, detail="Only list owner or admin can edit this list")
    if title is not None:
        lst.title = title
    if notebox is not None:
        lst.notebox = notebox
    if visibility is not None:
        lst.visibility = normalize_list_visibility(visibility, default=normalize_list_visibility(lst.visibility))
    db.commit()
    db.refresh(lst)
    list_asset_ids = [
        row[0]
        for row in db.query(PhotoListItem.asset_id).filter(PhotoListItem.list_id == lst.id).all()
        if row[0] is not None
    ]
    for item_asset_id in list_asset_ids:
        _refresh_asset_text_index_for_asset(db, tenant, item_asset_id)

    # Get creator display name
    created_by_name = None
    if lst.created_by_uid:
        user = db.query(UserProfile.display_name).filter(
            UserProfile.supabase_uid == lst.created_by_uid
        ).first()
        created_by_name = user[0] if user else None

    return _serialize_list_row(
        list_row=lst,
        created_by_name=created_by_name,
        item_count=len(lst.items),
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )


@router.delete("/{list_id:int}", response_model=dict)
async def delete_list(
    list_id: int,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Delete a list and its items."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    lst = _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if not can_edit_list(lst, user=current_user, is_tenant_admin=is_tenant_admin):
        raise HTTPException(status_code=403, detail="Only list owner or admin can delete this list")
    list_asset_ids = [
        row[0]
        for row in db.query(PhotoListItem.asset_id).filter(PhotoListItem.list_id == lst.id).all()
        if row[0] is not None
    ]
    db.delete(lst)
    db.commit()
    for item_asset_id in list_asset_ids:
        _refresh_asset_text_index_for_asset(db, tenant, item_asset_id)
    return {"deleted": True}


@router.get("/{list_id:int}/items", response_model=list)
async def get_list_items(
    list_id: int,
    ids_only: bool = False,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Get all items in a list, ordered by sort_order ascending."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if ids_only:
        items = (
            db.query(
                PhotoListItem.id,
                ImageMetadata.id.label("photo_id"),
                PhotoListItem.asset_id,
                PhotoListItem.sort_order,
                PhotoListItem.added_at,
            )
            .join(PhotoList, PhotoListItem.list_id == PhotoList.id)
            .outerjoin(
                ImageMetadata,
                and_(
                    PhotoListItem.asset_id == ImageMetadata.asset_id,
                    tenant_column_filter(ImageMetadata, tenant),
                ),
            )
            .filter(PhotoListItem.list_id == list_id, tenant_column_filter(PhotoList, tenant))
            .order_by(PhotoListItem.sort_order.asc(), PhotoListItem.added_at.asc(), PhotoListItem.id.asc())
            .all()
        )
        return [
            {
                "id": item_id,
                "photo_id": photo_id,
                "asset_id": str(asset_id) if asset_id else None,
                "sort_order": sort_order,
                "added_at": added_at,
            }
            for item_id, photo_id, asset_id, sort_order, added_at in items
        ]
    items = (
        db.query(PhotoListItem, ImageMetadata)
        .join(
            ImageMetadata,
            and_(
                PhotoListItem.asset_id == ImageMetadata.asset_id,
                tenant_column_filter(ImageMetadata, tenant),
            ),
        )
        .filter(
            PhotoListItem.list_id == list_id
        )
        .order_by(PhotoListItem.sort_order.asc(), PhotoListItem.added_at.asc(), PhotoListItem.id.asc())
        .all()
    )
    assets_by_id = load_assets_for_images(db, [img for _, img in items])
    image_ids = [img.id for _, img in items]
    asset_id_to_image_id = {img.asset_id: img.id for _, img in items if img.asset_id is not None}
    asset_ids = list(asset_id_to_image_id.keys())
    active_tag_type = get_tenant_setting(db, tenant.id, 'active_machine_tag_type', default='siglip')
    tags = db.query(MachineTag).filter(
        MachineTag.asset_id.in_(asset_ids),
        tenant_column_filter(MachineTag, tenant),
        MachineTag.tag_type == active_tag_type
    ).all() if asset_ids else []
    permatags = db.query(Permatag).filter(
        Permatag.asset_id.in_(asset_ids),
        tenant_column_filter(Permatag, tenant)
    ).all() if asset_ids else []
    variant_count_by_asset = {
        asset_id: int(count or 0)
        for asset_id, count in (
            db.query(
                AssetDerivative.asset_id,
                func.count(AssetDerivative.id),
            ).filter(
                AssetDerivative.asset_id.in_(asset_ids),
                AssetDerivative.deleted_at.is_(None),
            ).group_by(AssetDerivative.asset_id).all() if asset_ids else []
        )
    }

    # Load all keywords
    keyword_ids = set()
    for tag in tags:
        keyword_ids.add(tag.keyword_id)
    for permatag in permatags:
        keyword_ids.add(permatag.keyword_id)

    # Build keyword lookup map
    keywords_map = {}
    if keyword_ids:
        keywords_data = db.query(
            Keyword.id,
            Keyword.keyword,
            KeywordCategory.name
        ).join(
            KeywordCategory, Keyword.category_id == KeywordCategory.id
        ).filter(
            Keyword.id.in_(keyword_ids)
        ).all()
        for kw_id, kw_name, cat_name in keywords_data:
            keywords_map[kw_id] = {"keyword": kw_name, "category": cat_name}

    tags_by_image = {}
    for tag in tags:
        image_id = asset_id_to_image_id.get(tag.asset_id)
        if image_id is None:
            continue
        kw_info = keywords_map.get(tag.keyword_id, {"keyword": "unknown", "category": "unknown"})
        tags_by_image.setdefault(image_id, []).append({
            "keyword": kw_info["keyword"],
            "category": kw_info["category"],
            "confidence": round(tag.confidence, 2)
        })
    permatags_by_image = {}
    for permatag in permatags:
        image_id = asset_id_to_image_id.get(permatag.asset_id)
        if image_id is None:
            continue
        kw_info = keywords_map.get(permatag.keyword_id, {"keyword": "unknown", "category": "unknown"})
        permatags_by_image.setdefault(image_id, []).append({
            "id": permatag.id,
            "keyword": kw_info["keyword"],
            "category": kw_info["category"],
            "signum": permatag.signum
        })
    response_items = []
    for item, img in items:
        storage_info = _resolve_storage_or_409(
            image=img,
            tenant=tenant,
            db=db,
            assets_by_id=assets_by_id,
        )
        machine_tags = sorted(tags_by_image.get(img.id, []), key=lambda x: x['confidence'], reverse=True)
        image_permatags = permatags_by_image.get(img.id, [])
        calculated_tags = calculate_tags(machine_tags, image_permatags)
        response_items.append({
            "id": item.id,
            "photo_id": img.id,
            "asset_id": str(item.asset_id) if item.asset_id else storage_info.asset_id,
            "sort_order": item.sort_order,
            "added_at": item.added_at,
            "image": {
                "id": img.id,
                "asset_id": storage_info.asset_id,
                "variant_count": int(variant_count_by_asset.get(img.asset_id, 0)),
                "has_variants": int(variant_count_by_asset.get(img.asset_id, 0)) > 0,
                "filename": img.filename,
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "file_size": img.file_size,
                "dropbox_path": storage_info.source_key,
                "source_provider": storage_info.source_provider,
                "source_key": storage_info.source_key,
                "source_rev": storage_info.source_rev,
                "source_url": _build_source_url(storage_info, tenant, img),
                "camera_make": img.camera_make,
                "camera_model": img.camera_model,
                "lens_model": img.lens_model,
                "iso": img.iso,
                "aperture": img.aperture,
                "shutter_speed": img.shutter_speed,
                "focal_length": img.focal_length,
                "gps_latitude": img.gps_latitude,
                "gps_longitude": img.gps_longitude,
                "capture_timestamp": img.capture_timestamp.isoformat() if img.capture_timestamp else None,
                "modified_time": img.modified_time.isoformat() if img.modified_time else None,
                "created_at": img.created_at.isoformat() if img.created_at else None,
                "thumbnail_path": storage_info.thumbnail_key,
                "thumbnail_url": storage_info.thumbnail_url,
                "tags_applied": img.tags_applied,
                "faces_detected": img.faces_detected,
                "rating": img.rating,
                "tags": machine_tags,
                "permatags": image_permatags,
                "calculated_tags": calculated_tags
            }
        })
    return response_items


@router.get("/{list_id:int}/export/pptx")
async def export_list_pptx(
    list_id: int,
    template_id: Optional[str] = Query(default=None),
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Generate and download a PPTX deck with full-res images from a list."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    list_row = _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )

    rows = (
        db.query(PhotoListItem, ImageMetadata)
        .join(
            ImageMetadata,
            and_(
                PhotoListItem.asset_id == ImageMetadata.asset_id,
                tenant_column_filter(ImageMetadata, tenant),
            ),
        )
        .filter(PhotoListItem.list_id == list_id)
        .order_by(PhotoListItem.sort_order.asc(), PhotoListItem.added_at.asc(), PhotoListItem.id.asc())
        .all()
    )

    if not rows:
        raise HTTPException(status_code=404, detail="This list has no images to export.")
    if len(rows) > PPTX_EXPORT_MAX_ITEMS:
        raise HTTPException(
            status_code=400,
            detail=(
                f"List has {len(rows)} items. On-demand PPTX export is limited to "
                f"{PPTX_EXPORT_MAX_ITEMS} items."
            ),
        )

    assets_by_id = load_assets_for_images(db, [image for _, image in rows])
    export_rows = []
    skipped_unavailable = 0
    for _, image in rows:
        try:
            storage_info = _resolve_storage_or_409(
                image=image,
                tenant=tenant,
                db=db,
                assets_by_id=assets_by_id,
                require_source=True,
            )
        except HTTPException:
            skipped_unavailable += 1
            continue

        provider_name, source_ref = _resolve_provider_ref(storage_info, image)
        if not source_ref:
            skipped_unavailable += 1
            continue
        export_rows.append(
            {
                "image_id": int(image.id),
                "filename": image.filename or f"image-{image.id}",
                "provider_name": provider_name,
                "source_ref": source_ref,
            }
        )

    if not export_rows:
        raise HTTPException(status_code=409, detail="No exportable images were found in this list.")

    template_temp_path = None
    if template_id:
        template_row = _get_template_or_404(db=db, tenant=tenant, template_id=template_id)
        if not _template_visible_to_user(
            template=template_row,
            current_user=current_user,
            is_tenant_admin=is_tenant_admin,
        ):
            raise HTTPException(status_code=404, detail="Template not found")
        template_storage_key = str(template_row.storage_key or "").strip()
        if not template_storage_key:
            raise HTTPException(status_code=409, detail="Template file is not available")
        try:
            bucket = storage.Client(project=settings.gcp_project_id).bucket(tenant.get_storage_bucket(settings))
            template_blob = bucket.blob(template_storage_key)
            if not template_blob.exists():
                raise HTTPException(status_code=404, detail="Template file not found")
            template_bytes = template_blob.download_as_bytes()
            template_file = tempfile.NamedTemporaryFile(prefix="zoltag-pptx-template-", suffix=".pptx", delete=False)
            template_file.write(template_bytes)
            template_file.flush()
            template_file.close()
            template_temp_path = template_file.name
        except HTTPException:
            raise
        except Exception as exc:  # noqa: BLE001
            raise HTTPException(status_code=500, detail=f"Failed to load template: {exc}")

    try:
        temp_path, added_count, skipped_count = await run_in_threadpool(
            _build_list_pptx_file,
            export_rows=export_rows,
            tenant=tenant,
            template_path=template_temp_path,
        )
    except RuntimeError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    except Exception as exc:  # noqa: BLE001
        logger.exception("Failed to build PPTX export for list %s", list_id)
        raise HTTPException(status_code=500, detail=f"Failed to export PPTX: {exc}")
    finally:
        if template_temp_path:
            _cleanup_temp_file(template_temp_path)

    timestamp = datetime.utcnow().strftime("%Y%m%d-%H%M")
    filename = _sanitize_export_filename(
        list_row.title,
        f"-{timestamp}.pptx",
    )
    headers = {
        "X-Zoltag-Export-Images": str(added_count),
        "X-Zoltag-Export-Skipped": str(skipped_unavailable + skipped_count),
    }
    return FileResponse(
        path=temp_path,
        media_type=PPTX_MIME_TYPE,
        filename=filename,
        background=BackgroundTask(_cleanup_temp_file, temp_path),
        headers=headers,
    )


@router.patch("/{list_id:int}/items/reorder", response_model=dict)
async def reorder_list_items(
    list_id: int,
    req: ReorderListItemsRequest,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Reorder list items by item id for a single list."""
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    list_row = _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if not can_edit_list(list_row, user=current_user, is_tenant_admin=is_tenant_admin):
        raise HTTPException(status_code=403, detail="Only list owner or admin can reorder this list")

    requested_ids = [int(item_id) for item_id in (req.item_ids or []) if int(item_id) > 0]
    if not requested_ids:
        raise HTTPException(status_code=400, detail="item_ids is required")

    existing_items = (
        db.query(PhotoListItem.id)
        .filter(PhotoListItem.list_id == list_id)
        .order_by(PhotoListItem.sort_order.asc(), PhotoListItem.added_at.asc(), PhotoListItem.id.asc())
        .all()
    )
    existing_ids = [int(row[0]) for row in existing_items]
    if not existing_ids:
        return {"list_id": list_id, "updated": 0, "item_ids": []}

    existing_set = set(existing_ids)
    ordered = []
    seen = set()
    for item_id in requested_ids:
        if item_id in existing_set and item_id not in seen:
            ordered.append(item_id)
            seen.add(item_id)
    for item_id in existing_ids:
        if item_id not in seen:
            ordered.append(item_id)

    for idx, item_id in enumerate(ordered, start=1):
        db.query(PhotoListItem).filter(
            PhotoListItem.id == item_id,
            PhotoListItem.list_id == list_id,
        ).update({"sort_order": idx}, synchronize_session=False)
    db.commit()
    return {"list_id": list_id, "updated": len(ordered), "item_ids": ordered}


@router.post("/{list_id:int}/add-photo", response_model=dict)
async def add_photo_to_specific_list(
    list_id: int,
    req: AddPhotoRequest,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Add a photo to a specific list."""
    photo_id = req.photo_id

    # Verify list belongs to tenant
    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    _get_accessible_list_or_404(
        db=db,
        tenant=tenant,
        list_id=list_id,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )

    image = db.query(ImageMetadata.id, ImageMetadata.asset_id).filter(
        ImageMetadata.id == photo_id,
        tenant_column_filter(ImageMetadata, tenant)
    ).first()
    if not image:
        raise HTTPException(status_code=404, detail="Image not found")
    asset_id = image.asset_id
    if asset_id is None:
        raise HTTPException(status_code=409, detail="Image has no asset_id")

    # Prevent duplicate items in the list
    existing_item = db.query(PhotoListItem).filter(
        PhotoListItem.list_id == list_id,
        PhotoListItem.asset_id == asset_id
    ).first()
    if existing_item:
        # Silent no-op if already present
        return {
            "list_id": list_id,
            "item_id": existing_item.id,
            "photo_id": image.id,
            "asset_id": str(existing_item.asset_id) if existing_item.asset_id else None,
            "added_at": existing_item.added_at,
        }

    max_sort_order = (
        db.query(func.max(PhotoListItem.sort_order))
        .filter(PhotoListItem.list_id == list_id)
        .scalar()
    )
    next_sort_order = int(max_sort_order or 0) + 1
    item = PhotoListItem(list_id=list_id, asset_id=asset_id, sort_order=next_sort_order)
    db.add(item)
    db.commit()
    db.refresh(item)
    _refresh_asset_text_index_for_asset(db, tenant, asset_id)
    return {
        "list_id": list_id,
        "item_id": item.id,
        "photo_id": image.id,
        "asset_id": str(item.asset_id) if item.asset_id else None,
        "added_at": item.added_at,
    }


@router.post("/add-photo", response_model=dict)
async def add_photo_to_list(
    req: AddPhotoRequest,
    tenant: Tenant = Depends(get_tenant),
    db: Session = Depends(get_db),
    current_user: UserProfile = Depends(get_current_user),
):
    """Add a photo to the most recently created list. If no lists exist, create one and add the photo."""
    photo_id = req.photo_id
    image = db.query(ImageMetadata.id, ImageMetadata.asset_id).filter(
        ImageMetadata.id == photo_id,
        tenant_column_filter(ImageMetadata, tenant)
    ).first()
    if not image:
        raise HTTPException(status_code=404, detail="Image not found")
    asset_id = image.asset_id
    if asset_id is None:
        raise HTTPException(status_code=409, detail="Image has no asset_id")

    is_tenant_admin = is_tenant_admin_user(db, tenant, current_user)
    recent = get_most_recent_list(
        db,
        tenant,
        current_user=current_user,
        is_tenant_admin=is_tenant_admin,
    )
    if not recent:
        # Auto-create new list
        recent = assign_tenant_scope(PhotoList(
            title="Untitled List",
            visibility="shared",
            created_by_uid=current_user.supabase_uid if current_user else None,
        ), tenant)
        db.add(recent)
        db.commit()
        db.refresh(recent)
    # Prevent duplicate items in the list
    existing_item = db.query(PhotoListItem).filter(
        PhotoListItem.list_id == recent.id,
        PhotoListItem.asset_id == asset_id
    ).first()
    if existing_item:
        # Silent no-op if already present
        return {
            "list_id": recent.id,
            "item_id": existing_item.id,
            "photo_id": image.id,
            "asset_id": str(existing_item.asset_id) if existing_item.asset_id else None,
            "added_at": existing_item.added_at,
        }
    max_sort_order = (
        db.query(func.max(PhotoListItem.sort_order))
        .filter(PhotoListItem.list_id == recent.id)
        .scalar()
    )
    next_sort_order = int(max_sort_order or 0) + 1
    item = PhotoListItem(list_id=recent.id, asset_id=asset_id, sort_order=next_sort_order)
    db.add(item)
    db.commit()
    db.refresh(item)
    _refresh_asset_text_index_for_asset(db, tenant, asset_id)
    return {
        "list_id": recent.id,
        "item_id": item.id,
        "photo_id": image.id,
        "asset_id": str(item.asset_id) if item.asset_id else None,
        "added_at": item.added_at,
    }
